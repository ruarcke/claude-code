#!/usr/bin/env python3
"""
RCPCC Presentation Generator
Reads an Excel spreadsheet and generates a professional .pptx presentation.
Usage: python rcpcc_generator.py "path/to/spreadsheet.xlsx"
"""

import sys
import os
import re
import locale
from datetime import datetime

# Fix Windows console encoding
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

try:
    from openpyxl import load_workbook
except ImportError:
    print("ERRO: openpyxl não instalado. Execute: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERRO: python-pptx não instalado. Execute: pip install python-pptx")
    sys.exit(1)


# ============================================================
# COLORS
# ============================================================
BLACK = RGBColor(0x1a, 0x1a, 0x1a)
DARK = RGBColor(0x33, 0x33, 0x33)
MED = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x2C, 0x3E, 0x6B)
GOLD = RGBColor(0xC9, 0x9A, 0x2E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
LRED = RGBColor(0xFC, 0xE4, 0xE4)
LGREEN = RGBColor(0xE4, 0xFC, 0xE4)
LBLUE = RGBColor(0xE8, 0xF0, 0xFF)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)

I = Inches


# ============================================================
# HELPER FUNCTIONS
# ============================================================
def shp(s, l, t, w, h, c=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if c:
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
    else:
        sh.fill.background()
    return sh


def rnd(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    return sh


def txt(s, l, t, w, h, text, sz=18, c=BLACK, b=False, a=PP_ALIGN.LEFT, f='Calibri'):
    tx = s.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = f
    p.alignment = a
    return tx


def line(s, l, t, w, c=BLACK, th=2):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(th))
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()


def circ(s, l, t, sz, c):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh


def top(s):
    shp(s, I(0), I(0), I(13.333), I(0.08), ACCENT)


def bot(s):
    shp(s, I(0), I(7.0), I(13.333), I(0.5), ACCENT)
    txt(s, I(0.8), I(7.05), I(5), I(0.4),
        "RCK Advogados  |  OAB/SP 405.599", sz=10, c=WHITE, b=True)
    txt(s, I(8), I(7.05), I(5), I(0.4),
        "Documento Confidencial", sz=10, c=RGBColor(0x88, 0x88, 0x99), a=PP_ALIGN.RIGHT)


def fmt_brl(val):
    """Format number as R$ X.XXX,XX"""
    if val is None or val == 0:
        return "R$ 0,00"
    try:
        v = float(val)
        neg = v < 0
        v = abs(v)
        intpart = int(v)
        dec = round((v - intpart) * 100)
        if dec >= 100:
            intpart += 1
            dec -= 100
        # Format with dots for thousands
        s = f"{intpart:,}".replace(",", ".")
        result = f"R$ {s},{dec:02d}"
        if neg:
            result = f"-{result}"
        return result
    except (ValueError, TypeError):
        return "R$ 0,00"


def fmt_pct(val):
    """Format as percentage"""
    if val is None:
        return "0%"
    try:
        v = float(val)
        if v < 1:
            v *= 100
        return f"{v:.0f}%"
    except (ValueError, TypeError):
        return "0%"


def data_extenso():
    """Return current date in Portuguese"""
    meses = [
        "", "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
        "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"
    ]
    d = datetime.now()
    return f"{d.day} de {meses[d.month]} de {d.year}"


# ============================================================
# READ SPREADSHEET
# ============================================================
def read_spreadsheet(path):
    """Read RCPCC spreadsheet and return structured data."""
    wb = load_workbook(path, data_only=True)
    ws = wb.active

    # Extract client name from filename
    fname = os.path.basename(path)
    # Pattern: "Tabela atendimento - RCPCC - CLIENT_NAME.xlsx"
    match = re.search(r'RCPCC\s*-\s*(.+?)\.xlsx', fname, re.IGNORECASE)
    if match:
        client_name = match.group(1).strip()
        # Clean up extra info like "- reuniao 7-8-25"
        client_name = re.sub(r'\s*-\s*reuniao?\s.*$', '', client_name, flags=re.IGNORECASE)
        # Clean phone numbers
        client_name = re.sub(r'^\+?\d[\d\s-]+$', 'Cliente', client_name)
    else:
        client_name = "Cliente"

    # Read debt rows (A2:G14, filter out placeholders)
    debts = []
    for row in range(2, 15):
        banco = ws[f'A{row}'].value
        if banco and str(banco).strip().upper() != 'NOME' and banco != '':
            modalidade = ws[f'B{row}'].value or ''
            vl_parcela = ws[f'C{row}'].value or 0
            qtd = ws[f'D{row}'].value or 0
            pagas = ws[f'E{row}'].value or 0
            atraso = ws[f'F{row}'].value or 0
            total = ws[f'G{row}'].value or 0

            if float(vl_parcela) > 0 or float(total) > 0:
                debts.append({
                    'banco': str(banco).strip(),
                    'modalidade': str(modalidade).strip(),
                    'vl_parcela': float(vl_parcela),
                    'qtd': int(float(qtd)),
                    'pagas': int(float(pagas)),
                    'atraso': int(float(atraso)),
                    'total': float(total),
                })

    # Fixed cells
    total_geral = float(ws['G15'].value or 0)
    pagamento = float(ws['G16'].value or 0)
    pct_desconto_raw = ws['H16'].value
    economia = float(ws['G17'].value or 0)
    honorarios = float(ws['G18'].value or 0)
    parcelas_hon = int(float(ws['F19'].value or 12))
    vl_parcela_hon = float(ws['G19'].value or 0)

    # Calculate REAL discount percentage from actual values
    # H16 is the configured target but cooperatives get less,
    # so we calculate from actual total vs payment
    if total_geral > 0 and pagamento > 0:
        pct_desc = (1 - pagamento / total_geral) * 100
    elif pct_desconto_raw:
        pct_desc = float(pct_desconto_raw)
        if pct_desc < 1:
            pct_desc *= 100
    else:
        pct_desc = 85

    # Search for "à vista" and "cartão" values (positions vary)
    avista = None
    cartao = None
    for row in range(10, 30):
        k_val = ws[f'K{row}'].value
        l_val = ws[f'L{row}'].value
        if k_val and 'vista' in str(k_val).lower() and l_val:
            avista = float(l_val)
        if k_val and ('cart' in str(k_val).lower() or 'cartão' in str(k_val).lower()) and l_val:
            cartao = float(l_val)

    # Fallback calculations
    if avista is None:
        avista = honorarios * 0.9
    if cartao is None:
        cartao = honorarios * 0.95

    # Count unique banks and total debts
    bancos_unicos = len(set(d['banco'] for d in debts))
    total_dividas = len(debts)
    todas_atraso = all(d['atraso'] > 0 for d in debts)

    data = {
        'client_name': client_name,
        'debts': debts,
        'total_geral': total_geral,
        'pagamento': pagamento,
        'pct_desconto': pct_desc,
        'economia': economia,
        'honorarios': honorarios,
        'parcelas_hon': parcelas_hon,
        'vl_parcela_hon': vl_parcela_hon,
        'avista': avista,
        'cartao': cartao,
        'bancos_unicos': bancos_unicos,
        'total_dividas': total_dividas,
        'todas_atraso': todas_atraso,
    }

    wb.close()
    return data


# ============================================================
# GENERATE PRESENTATION
# ============================================================
def generate_presentation(data, output_path):
    prs = Presentation()
    prs.slide_width = I(13.333)
    prs.slide_height = I(7.5)

    nome = data['client_name']
    primeiro_nome = nome.split()[0] if nome else "Cliente"

    # ============================================================
    # SLIDE 1 - CAPA
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(4.5), I(1.0), I(4.333), I(1.5), "RCK", sz=80, c=BLACK, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(4.5), I(2.4), I(4.333), I(0.6), "A D V O G A D O S", sz=20, c=MED,
        a=PP_ALIGN.CENTER, f='Calibri Light')
    line(sl, I(5.5), I(3.3), I(2.333), GOLD, 2)
    txt(sl, I(2), I(3.7), I(9.333), I(0.6), "DIAGNÓSTICO FINANCEIRO PERSONALIZADO",
        sz=22, c=DARK, a=PP_ALIGN.CENTER, f='Calibri Light')
    txt(sl, I(2), I(4.5), I(9.333), I(0.8), nome,
        sz=40, c=BLACK, b=True, a=PP_ALIGN.CENTER)
    line(sl, I(5.5), I(5.5), I(2.333), GOLD, 2)
    txt(sl, I(2), I(5.8), I(9.333), I(0.4),
        "Preparado por Dr. Ruarcke Oliveira  |  OAB/SP 405.599", sz=14, c=MED, a=PP_ALIGN.CENTER)
    txt(sl, I(2), I(6.2), I(9.333), I(0.4), data_extenso(), sz=13, c=MED, a=PP_ALIGN.CENTER)
    shp(sl, I(0), I(7.0), I(13.333), I(0.5), ACCENT)
    txt(sl, I(0), I(7.05), I(13.333), I(0.4),
        f"CONFIDENCIAL — Preparado exclusivamente para {nome}", sz=11, c=WHITE, a=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 2 - SEU ADVOGADO
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(8), I(0.6), "SEU ADVOGADO", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)
    txt(sl, I(0.8), I(1.4), I(6), I(0.6), "Dr. Ruarcke Antônio Diniz de Oliveira", sz=24, c=BLACK, b=True)
    txt(sl, I(0.8), I(1.95), I(8), I(0.4),
        "Advogado especializado em Direito Processual Civil e Direito Bancário", sz=14, c=MED)

    stats = [
        ("5.773+", "Processos\natuados", ACCENT),
        ("6", "Inscrições\nna OAB", BLUE),
        ("6", "Estados de\natuação", GREEN),
        ("R$50M+", "Economia gerada\npara clientes", GOLD),
    ]
    for i, (n, d, c) in enumerate(stats):
        x = I(0.8 + i * 3.1)
        rnd(sl, x, I(2.6), I(2.8), I(1.8), c)
        txt(sl, x, I(2.8), I(2.8), I(0.8), n, sz=36, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x, I(3.5), I(2.8), I(0.7), d, sz=13, c=WHITE, a=PP_ALIGN.CENTER)

    txt(sl, I(0.8), I(4.7), I(5), I(0.4), "Inscrições na Ordem dos Advogados do Brasil", sz=16, c=BLACK, b=True)
    oabs = [
        ("SP", "405.599", "Principal"), ("BA", "89.746", "Suplementar"),
        ("CE", "55.655-A", "Suplementar"), ("MG", "237.220", "Suplementar"),
        ("PE", "68.097", "Suplementar"), ("SC", "77.135-A", "Suplementar"),
    ]
    for i, (uf, num, tipo) in enumerate(oabs):
        x = I(0.8 + i * 2.05)
        y = I(5.2)
        c2 = ACCENT if tipo == "Principal" else BLUE
        rnd(sl, x, y, I(1.85), I(1.1), c2)
        txt(sl, x, y + I(0.1), I(1.85), I(0.4), f"OAB/{uf}", sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x, y + I(0.5), I(1.85), I(0.3), num, sz=13,
            c=RGBColor(0xCC, 0xCC, 0xDD), a=PP_ALIGN.CENTER)
        txt(sl, x, y + I(0.75), I(1.85), I(0.25), tipo, sz=10,
            c=RGBColor(0xAA, 0xAA, 0xBB), a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 3 - SOBRE O ESCRITÓRIO
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(8), I(0.6), "SOBRE O ESCRITÓRIO", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)
    txt(sl, I(0.8), I(1.4), I(5.5), I(2.5),
        "O escritório RCK Advogados é referência em soluções jurídicas para renegociação de dívidas "
        "bancárias. Com atuação em 6 estados brasileiros, nossa equipe busca os melhores descontos e "
        "condições para devolver a tranquilidade financeira aos nossos clientes.", sz=15, c=MED)
    txt(sl, I(7), I(1.4), I(5.5), I(0.4), "Nossas Especialidades", sz=18, c=BLACK, b=True)
    svcs = [
        "Renegociação de dívidas bancárias (descontos de até 90%)",
        "Renegociação com cooperativas (descontos de 50% a 70%)",
        "Plano de pagamento personalizado",
        "Provisionamento bancário estratégico",
        "Defesa em processos de execução",
        "Defesa em ações monitórias",
        "Análise de contratos e identificação de irregularidades",
        "Monitoramento processual contínuo",
    ]
    for i, s in enumerate(svcs):
        y = I(1.9 + i * 0.5)
        rnd(sl, I(7), y, I(0.3), I(0.3), GREEN)
        txt(sl, I(7.05), y - I(0.02), I(0.3), I(0.3), "✓", sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, I(7.5), y, I(5), I(0.4), s, sz=13, c=DARK)
    cards = [
        ("90%", "Desconto máximo\nem bancos", GREEN),
        ("70%", "Desconto máximo\nem cooperativas", BLUE),
        ("500+", "Casos de RCPCC\nresolvidos", ACCENT),
    ]
    for i, (n, d, c) in enumerate(cards):
        x = I(0.8 + i * 2.1)
        rnd(sl, x, I(4.5), I(1.9), I(2.0), c)
        txt(sl, x, I(4.8), I(1.9), I(0.7), n, sz=32, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x, I(5.5), I(1.9), I(0.7), d, sz=12, c=WHITE, a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 4 - DIAGNÓSTICO (from spreadsheet)
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(8), I(0.6), "DIAGNÓSTICO FINANCEIRO", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)
    txt(sl, I(0.8), I(1.25), I(8), I(0.4),
        f"Entendemos a sua situação, {primeiro_nome}. Veja o panorama completo:", sz=14, c=MED)
    rnd(sl, I(8.5), I(0.4), I(4), I(0.5), ORANGE)
    txt(sl, I(8.5), I(0.43), I(4), I(0.4),
        "⚡ EDITÁVEL — Atualize valores na reunião", sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    # Table header
    hy = I(1.7)
    shp(sl, I(0.8), hy, I(11.733), I(0.5), ACCENT)
    hdrs = [
        ("BANCO", 0.8, 2.5), ("MODALIDADE", 3.3, 2.2), ("VL. PARCELA", 5.5, 1.5),
        ("QTD.", 7.0, 0.7), ("PAGAS", 7.7, 0.7), ("ATRASO", 8.4, 0.7),
        ("TOTAL", 9.1, 1.8), ("% TOTAL", 10.9, 1.6),
    ]
    for t, x, w in hdrs:
        txt(sl, I(x), hy + I(0.05), I(w), I(0.4), t, sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    # Data rows (real debts + placeholder rows up to 8)
    total_g = data['total_geral']
    rows_data = []
    for d in data['debts']:
        pct = f"{d['total'] / total_g * 100:.1f}%" if total_g > 0 else "0%"
        rows_data.append((
            d['banco'], d['modalidade'], fmt_brl(d['vl_parcela']),
            str(d['qtd']), str(d['pagas']), str(d['atraso']),
            fmt_brl(d['total']), pct
        ))

    # Fill remaining with editable placeholders
    while len(rows_data) < 8:
        rows_data.append(("[Banco]", "[Tipo]", "R$ 0", "0", "0", "0", "R$ 0", "0%"))

    for i, row in enumerate(rows_data):
        ry = I(2.2 + i * 0.45)
        is_placeholder = row[0].startswith("[")
        rc = RGBColor(0xFF, 0xFB, 0xE6) if is_placeholder else (LGRAY if i % 2 == 0 else WHITE)
        shp(sl, I(0.8), ry, I(11.733), I(0.45), rc)
        cols = [(0.8, 2.5), (3.3, 2.2), (5.5, 1.5), (7.0, 0.7),
                (7.7, 0.7), (8.4, 0.7), (9.1, 1.8), (10.9, 1.6)]
        for j, (x, w) in enumerate(cols):
            cc = RED if j == 5 and row[5] not in ("0", "0%") else (MED if is_placeholder else DARK)
            txt(sl, I(x), ry + I(0.03), I(w), I(0.35), row[j], sz=11, c=cc, b=(j == 0), a=PP_ALIGN.CENTER)

    # Total row
    ty = I(5.8)
    shp(sl, I(0.8), ty, I(11.733), I(0.5), RED)
    txt(sl, I(0.8), ty + I(0.05), I(8.3), I(0.4), "TOTAL DAS DÍVIDAS",
        sz=14, c=WHITE, b=True, a=PP_ALIGN.RIGHT)
    txt(sl, I(9.1), ty + I(0.05), I(1.8), I(0.4), fmt_brl(total_g),
        sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    # Status badges
    if data['todas_atraso']:
        rnd(sl, I(0.8), I(6.4), I(5.5), I(0.5), LRED)
        txt(sl, I(1.0), I(6.42), I(5), I(0.4),
            "⚠ Todas as parcelas em atraso — Situação crítica", sz=12, c=RED, b=True)
    rnd(sl, I(7), I(6.4), I(5.5), I(0.5), RGBColor(0xFF, 0xFB, 0xE6))
    txt(sl, I(7.2), I(6.42), I(5), I(0.4),
        "💡 Linhas amarelas: edite durante a reunião", sz=11, c=ORANGE, b=True)
    bot(sl)

    # ============================================================
    # SLIDE 5 - CONSEQUÊNCIAS
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(10), I(0.6), "O QUE ACONTECE SE VOCÊ NÃO AGIR?",
        sz=28, c=RED, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), RED, 3)
    txt(sl, I(0.8), I(1.25), I(10), I(0.4),
        "As consequências de não resolver suas dívidas bancárias podem ser graves e irreversíveis", sz=15, c=MED)

    consequences = [
        ("⚖️", "AÇÃO MONITÓRIA",
         "O banco entra com ação monitória. Até gerar penhora dentro do processo leva ~18 meses. "
         "Você será citado e terá custos processuais.", "Penhora: ~18 meses", RED),
        ("🔨", "EXECUÇÃO DIRETA",
         "Empréstimo pessoal permite execução direta — sem monitória prévia. O juiz pode determinar "
         "penhora dos seus bens e bloqueio de contas imediatamente.", "Risco: Imediato", DARK_RED),
        ("🏦", "BLOQUEIO DE CONTAS",
         "Bloqueio judicial (penhora online) de todas as suas contas bancárias via BacenJud. "
         "Salário é protegido por lei, mas dinheiro em conta pode ser penhorado.", "Impacto: Imediato",
         RGBColor(0x8B, 0x45, 0x13)),
        ("📉", "NOME NEGATIVADO",
         "CPF nos órgãos de proteção (SPC/Serasa), impedindo financiamentos, cartões e até aluguéis. "
         "Após acordo, limite bancário reduzido por ~12 meses.", "Duração: Até 5 anos", ORANGE),
        ("💰", "JUROS CRESCENTES",
         "Juros e multas continuam incidindo. A dívida pode chegar a 180-200% do valor original, "
         "tornando a quitação cada vez mais difícil.", "Crescimento: Exponencial",
         RGBColor(0x8B, 0x00, 0x8B)),
        ("🚗", "PENHORA DE BENS",
         "Veículos e dinheiro em conta podem ser penhorados. Imóvel de moradia e salário são "
         "protegidos por lei. Dica: transferir veículos antes da ação.", "Proteção: Limitada", BLACK),
    ]
    for i, (icon, title, desc, tag, color) in enumerate(consequences):
        col = i % 3
        row = i // 3
        x = I(0.5 + col * 4.2)
        y = I(1.7 + row * 2.7)
        rnd(sl, x, y, I(4.0), I(2.5), LGRAY)
        circ(sl, x + I(0.15), y + I(0.15), I(0.55), color)
        txt(sl, x + I(0.15), y + I(0.17), I(0.55), I(0.5), icon, sz=18, c=WHITE, a=PP_ALIGN.CENTER)
        txt(sl, x + I(0.85), y + I(0.2), I(2.9), I(0.4), title, sz=15, c=BLACK, b=True)
        rnd(sl, x + I(0.85), y + I(0.55), I(1.8), I(0.25), color)
        txt(sl, x + I(0.85), y + I(0.56), I(1.8), I(0.22), tag, sz=9, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x + I(0.15), y + I(0.95), I(3.6), I(1.3), desc, sz=11, c=MED)
    bot(sl)

    # ============================================================
    # SLIDE 6 - SOMOS A SOLUÇÃO
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(10), I(0.6), "MAS EXISTE UMA SOLUÇÃO", sz=28, c=GREEN, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GREEN, 3)
    txt(sl, I(0.8), I(1.25), I(10), I(0.4),
        "A RCK Advogados atua estrategicamente para resolver suas dívidas com os maiores descontos do mercado",
        sz=15, c=MED)

    steps = [
        ("01", "Análise Completa",
         "Estudamos cada dívida,\nidentificamos irregularidades\ne traçamos a estratégia ideal", GREEN),
        ("02", "Negociação Técnica",
         "Entramos em contato com\ncada banco usando técnicas\njurídicas para máximo desconto", BLUE),
        ("03", "Acordo Favorável",
         "Formalizamos acordos com\ndescontos de 80-90% (bancos)\nou 50-70% (cooperativas)", ACCENT),
        ("04", "Acompanhamento",
         "Monitoramos seu CPF,\ndefendemos em execuções\ne garantimos sua proteção", GOLD),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        x = I(0.5 + i * 3.2)
        y = I(1.8)
        rnd(sl, x, y, I(3.0), I(3.2), LGRAY)
        circ(sl, x + I(1.0), y + I(0.2), I(0.9), color)
        txt(sl, x + I(1.0), y + I(0.32), I(0.9), I(0.6), num, sz=26, c=WHITE, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x + I(0.15), y + I(1.3), I(2.7), I(0.4), title, sz=16, c=BLACK, b=True, a=PP_ALIGN.CENTER)
        txt(sl, x + I(0.15), y + I(1.75), I(2.7), I(1.2), desc, sz=12, c=MED, a=PP_ALIGN.CENTER)

    rnd(sl, I(0.8), I(5.3), I(11.7), I(1.4), GREEN)
    txt(sl, I(0.8), I(5.4), I(11.7), I(0.5), "RESULTADO PARA VOCÊ",
        sz=20, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(5.9), I(3.9), I(0.6), "Descontos de até 90%",
        sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(4.7), I(5.9), I(3.9), I(0.6), "Nome limpo novamente",
        sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(8.6), I(5.9), I(3.9), I(0.6), "Tranquilidade financeira",
        sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 7 - PROVA SOCIAL
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(10), I(0.6), "RESULTADOS REAIS — CASO COMPROVADO NO TJSP",
        sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)
    txt(sl, I(0.8), I(1.25), I(10), I(0.4),
        "Processo público nº 1000361-89.2024 — Acordo homologado judicialmente", sz=14, c=MED)

    rnd(sl, I(0.8), I(1.8), I(11.7), I(3.0), LGRAY)
    rnd(sl, I(0.8), I(1.8), I(11.7), I(0.55), ACCENT)
    txt(sl, I(0.8), I(1.85), I(11.7), I(0.45),
        "CASO REAL — ACORDO HOMOLOGADO PELO TRIBUNAL DE JUSTIÇA DE SÃO PAULO",
        sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    rnd(sl, I(1.1), I(2.6), I(3.5), I(1.8), LRED)
    txt(sl, I(1.1), I(2.65), I(3.5), I(0.3), "DÍVIDA ORIGINAL", sz=12, c=RED, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(1.1), I(3.0), I(3.5), I(0.6), "R$ 267.412,36", sz=34, c=RED, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(1.1), I(3.65), I(3.5), I(0.3), "Banco credor em execução judicial",
        sz=11, c=RED, a=PP_ALIGN.CENTER)

    txt(sl, I(4.8), I(3.1), I(1.5), I(0.8), "→", sz=50, c=GREEN, b=True, a=PP_ALIGN.CENTER)

    rnd(sl, I(6.5), I(2.6), I(3.5), I(1.8), LGREEN)
    txt(sl, I(6.5), I(2.65), I(3.5), I(0.3), "VALOR DO ACORDO",
        sz=12, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(3.0), I(3.5), I(0.6), "R$ 26.741,24", sz=34, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(3.65), I(3.5), I(0.3), "Parcela única à vista — quitado",
        sz=11, c=GREEN, a=PP_ALIGN.CENTER)

    rnd(sl, I(10.3), I(2.6), I(2.0), I(1.8), GREEN)
    txt(sl, I(10.3), I(2.8), I(2.0), I(0.7), "90%", sz=44, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(10.3), I(3.5), I(2.0), I(0.5), "DESCONTO\nOBTIDO", sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    rnd(sl, I(0.8), I(4.9), I(11.7), I(0.75), RGBColor(0xFD, 0xF5, 0xE6))
    txt(sl, I(1.0), I(4.95), I(11.3), I(0.3),
        "📄 TRECHO DO ACORDO JUDICIAL HOMOLOGADO:", sz=12, c=DARK, b=True)
    txt(sl, I(1.0), I(5.25), I(11.3), I(0.3),
        "\"Valor do acordo: R$ 26.741,24 [...] Será concedido o desconto condicional no valor de R$ 240.671,12\"",
        sz=12, c=MED)

    rnd(sl, I(0.8), I(5.8), I(5.5), I(0.5), ORANGE)
    txt(sl, I(0.8), I(5.83), I(5.5), I(0.4),
        "💡 Vou abrir o processo real na tela para você verificar", sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(5.83), I(6), I(0.22), "🔗 Bradesco: esaj.tjsp.jus.br — Foro 224", sz=11, c=BLUE)
    txt(sl, I(6.5), I(6.05), I(6), I(0.22), "🔗 Santander: esaj.tjsp.jus.br — Foro 541", sz=11, c=BLUE)

    rnd(sl, I(0.8), I(6.4), I(11.7), I(0.5), ACCENT)
    txt(sl, I(0.8), I(6.43), I(11.7), I(0.4),
        "📊 5.773 processos no Escavador  |  4.649 em SP  |  Consulte: escavador.com/nomes/ruarcke-antonio-diniz-de-oliveira",
        sz=11, c=RGBColor(0xBB, 0xBB, 0xCC), a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 8 - ENTREGÁVEIS
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(10), I(0.6), "O QUE ESTÁ INCLUSO NO SEU PLANO", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)
    txt(sl, I(0.8), I(1.25), I(10), I(0.4),
        "Tudo que você recebe ao contratar a RCK Advogados", sz=14, c=MED)

    delivs = [
        ("📞", "2 Contatos Mensais",
         "Pelo menos 2 atualizações por mês sobre o andamento da negociação e o percentual de desconto obtido.",
         GREEN),
        ("🔍", "Monitoramento do CPF",
         "Acompanhamento contínuo para verificar se algum banco entrou com processo judicial contra você.",
         BLUE),
        ("📋", "Análise de Contratos",
         "Revisão da documentação para identificar irregularidades nos contratos bancários que podem ser usadas a seu favor.",
         ACCENT),
        ("🛡️", "Defesa em Execuções",
         "Se o banco entrar com processo de execução, sua defesa já está inclusa. Protegemos seus bens e contas.",
         RED),
        ("⚖️", "Impugnação e Desbloqueio",
         "Em caso de bloqueio de conta, pedido de desbloqueio em aproximadamente 7 dias úteis.",
         ORANGE),
        ("📊", "Negociação Estratégica",
         "Atuação técnica junto aos bancos por 12-16 meses para obter os maiores descontos possíveis (80-90%).",
         GOLD),
    ]
    for i, (icon, title, desc, color) in enumerate(delivs):
        col = i % 3
        row = i // 3
        x = I(0.6 + col * 4.15)
        y = I(1.7 + row * 2.6)
        rnd(sl, x, y, I(3.95), I(2.3), LGRAY)
        circ(sl, x + I(0.2), y + I(0.2), I(0.6), color)
        txt(sl, x + I(0.2), y + I(0.22), I(0.6), I(0.55), icon, sz=20, c=WHITE, a=PP_ALIGN.CENTER)
        txt(sl, x + I(1.0), y + I(0.25), I(2.7), I(0.4), title, sz=15, c=BLACK, b=True)
        txt(sl, x + I(0.2), y + I(0.9), I(3.5), I(1.2), desc, sz=12, c=MED)
    bot(sl)

    # ============================================================
    # SLIDE 9 - PROPOSTA (before vs after from spreadsheet)
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(10), I(0.6), "SUA PROPOSTA PERSONALIZADA", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)

    rnd(sl, I(0.8), I(1.3), I(11.5), I(0.45), LBLUE)
    txt(sl, I(1), I(1.32), I(11), I(0.4),
        "ℹ Bancos: descontos de até 85-90%  |  Cooperativas: descontos de 50-70%  |  Valores editáveis",
        sz=12, c=BLUE, b=True, a=PP_ALIGN.CENTER)

    # Before
    rnd(sl, I(0.8), I(2.0), I(3.8), I(4.0), LRED)
    rnd(sl, I(0.8), I(2.0), I(3.8), I(0.55), RED)
    txt(sl, I(0.8), I(2.05), I(3.8), I(0.45), "SITUAÇÃO ATUAL", sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(2.9), I(3.8), I(0.4), "Total das dívidas", sz=14, c=MED, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(3.3), I(3.8), I(0.7), fmt_brl(data['total_geral']),
        sz=36, c=RED, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(4.1), I(3.8), I(0.4),
        f"{data['total_dividas']} dívida{'s' if data['total_dividas'] != 1 else ''} "
        f"em {data['bancos_unicos']} banco{'s' if data['bancos_unicos'] != 1 else ''}",
        sz=13, c=MED, a=PP_ALIGN.CENTER)
    status_text = "Todas em atraso" if data['todas_atraso'] else "Parcelas em atraso"
    txt(sl, I(0.8), I(4.5), I(3.8), I(0.4), status_text, sz=13, c=RED, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(5.0), I(3.8), I(0.4), "Risco de execução judicial", sz=12, c=RED, a=PP_ALIGN.CENTER)

    txt(sl, I(4.8), I(3.5), I(1.5), I(1), "→", sz=60, c=GREEN, b=True, a=PP_ALIGN.CENTER)

    # After
    rnd(sl, I(6.5), I(2.0), I(3.8), I(4.0), LGREEN)
    rnd(sl, I(6.5), I(2.0), I(3.8), I(0.55), GREEN)
    txt(sl, I(6.5), I(2.05), I(3.8), I(0.45), "COM A RCK ADVOGADOS",
        sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(2.9), I(3.8), I(0.4), "Valor negociado", sz=14, c=MED, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(3.3), I(3.8), I(0.7), fmt_brl(data['pagamento']),
        sz=36, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(4.1), I(3.8), I(0.4), f"Desconto de {data['pct_desconto']:.0f}%",
        sz=13, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(4.5), I(3.8), I(0.4), "Dívidas quitadas", sz=13, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(6.5), I(5.0), I(3.8), I(0.4), "Nome limpo novamente", sz=12, c=GREEN, a=PP_ALIGN.CENTER)

    rnd(sl, I(2.5), I(6.2), I(8.333), I(0.6), GREEN)
    txt(sl, I(2.5), I(6.23), I(8.333), I(0.5),
        f"💰 ECONOMIA TOTAL:  {fmt_brl(data['economia'])}",
        sz=24, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 10 - CONDIÇÕES DE PAGAMENTO
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    top(sl)
    txt(sl, I(0.8), I(0.4), I(8), I(0.6), "CONDIÇÕES DE PAGAMENTO", sz=28, c=BLACK, b=True)
    line(sl, I(0.8), I(1.05), I(1.5), GOLD, 3)

    # Option 1 - À vista
    rnd(sl, I(0.8), I(1.4), I(5.5), I(4.0), LGRAY)
    rnd(sl, I(0.8), I(1.4), I(5.5), I(0.55), GREEN)
    txt(sl, I(0.8), I(1.45), I(5.5), I(0.45), "OPÇÃO 1 — À VISTA",
        sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(1.2), I(2.2), I(4.7), I(0.3), "Valor com 10% de desconto adicional", sz=13, c=MED)
    txt(sl, I(1.2), I(2.6), I(4.7), I(0.7), fmt_brl(data['avista']), sz=40, c=GREEN, b=True)
    txt(sl, I(1.2), I(3.4), I(4.7), I(0.3), "Pagamento único — maior economia", sz=13, c=MED)
    rnd(sl, I(1.2), I(3.9), I(4.7), I(0.4), GREEN)
    txt(sl, I(1.2), I(3.93), I(4.7), I(0.35), "✦ MELHOR CUSTO-BENEFÍCIO",
        sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    # Option 2 - Parcelado
    entrada = data['vl_parcela_hon'] * 3 if data['vl_parcela_hon'] > 0 else data['honorarios'] / 4
    parcela_val = data['vl_parcela_hon']
    parcelas_n = data['parcelas_hon']
    total_parcelado = entrada + (parcela_val * parcelas_n)

    rnd(sl, I(6.8), I(1.4), I(5.5), I(4.0), LGRAY)
    rnd(sl, I(6.8), I(1.4), I(5.5), I(0.55), BLUE)
    txt(sl, I(6.8), I(1.45), I(5.5), I(0.45), "OPÇÃO 2 — PARCELADO",
        sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(7.2), I(2.2), I(4.7), I(0.3),
        f"Entrada + {parcelas_n} parcelas mensais", sz=13, c=MED)
    txt(sl, I(7.2), I(2.6), I(2), I(0.3), "Entrada:", sz=13, c=MED)
    txt(sl, I(9.2), I(2.55), I(2.5), I(0.4), fmt_brl(entrada), sz=18, c=DARK, b=True)
    txt(sl, I(7.2), I(3.0), I(2), I(0.3), f"{parcelas_n}x de:", sz=13, c=MED)
    txt(sl, I(9.2), I(2.9), I(2.5), I(0.6), fmt_brl(parcela_val), sz=28, c=BLUE, b=True)
    txt(sl, I(7.2), I(3.5), I(4.7), I(0.3), f"Total: {fmt_brl(total_parcelado)}", sz=13, c=MED)
    rnd(sl, I(7.2), I(3.9), I(4.7), I(0.4), BLUE)
    txt(sl, I(7.2), I(3.93), I(4.7), I(0.35), "✦ MAIOR FLEXIBILIDADE",
        sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    # Honorários
    rnd(sl, I(0.8), I(5.7), I(11.533), I(1.0), LBLUE)
    txt(sl, I(0.8), I(5.75), I(11.533), I(0.4), "HONORÁRIOS DO ESCRITÓRIO",
        sz=14, c=BLUE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(0.8), I(6.15), I(11.533), I(0.4),
        f"Honorários iniciais: 10% da dívida  |  Honorários finais: 10% do desconto obtido  "
        f"|  Em até {parcelas_n}x de {fmt_brl(parcela_val)}",
        sz=13, c=BLUE, a=PP_ALIGN.CENTER)
    bot(sl)

    # ============================================================
    # SLIDE 11 - FECHAMENTO
    # ============================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = ACCENT

    txt(sl, I(4.5), I(0.8), I(4.333), I(1.2), "RCK", sz=60, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, I(4.5), I(1.7), I(4.333), I(0.5), "A D V O G A D O S",
        sz=18, c=RGBColor(0x99, 0x99, 0xAA), a=PP_ALIGN.CENTER, f='Calibri Light')
    line(sl, I(5.5), I(2.4), I(2.333), GOLD, 2)

    txt(sl, I(2), I(2.8), I(9.333), I(0.6),
        f"{primeiro_nome}, vamos resolver isso juntos?", sz=30, c=WHITE, b=True, a=PP_ALIGN.CENTER)

    rnd(sl, I(2.5), I(3.6), I(8.333), I(2.5), RGBColor(0x25, 0x25, 0x40))
    txt(sl, I(2.5), I(3.7), I(8.333), I(0.4), "RESUMO DA SUA SOLUÇÃO",
        sz=16, c=GOLD, b=True, a=PP_ALIGN.CENTER)

    summary = [
        ("Dívida atual:", fmt_brl(data['total_geral']), RED),
        ("Valor com desconto:", fmt_brl(data['pagamento']), GREEN),
        ("Sua economia:", f"{fmt_brl(data['economia'])} ({data['pct_desconto']:.0f}%)", GREEN),
        ("Honorários a partir de:", f"{parcelas_n}x de {fmt_brl(parcela_val)}", WHITE),
        ("Entregáveis:", "Negociação + defesa + monitoramento", WHITE),
    ]
    for i, (label, value, color) in enumerate(summary):
        y_pos = I(4.2 + i * 0.3)
        txt(sl, I(3), y_pos, I(3.5), I(0.3), label, sz=14, c=RGBColor(0xAA, 0xAA, 0xBB))
        txt(sl, I(6.5), y_pos, I(4), I(0.3), value, sz=14, c=color, b=True)

    line(sl, I(5.5), I(6.3), I(2.333), GOLD, 2)
    txt(sl, I(2), I(6.5), I(9.333), I(0.3), "Dr. Ruarcke Oliveira  |  OAB/SP 405.599",
        sz=13, c=RGBColor(0xAA, 0xAA, 0xBB), a=PP_ALIGN.CENTER)
    txt(sl, I(2), I(6.8), I(9.333), I(0.3),
        "contato@rckadvogados.com.br  |  www.rckadvogados.com.br",
        sz=12, c=RGBColor(0x88, 0x88, 0x99), a=PP_ALIGN.CENTER)
    txt(sl, I(2), I(7.1), I(9.333), I(0.3),
        f"Confidencial — Preparado exclusivamente para {nome}",
        sz=10, c=RGBColor(0x66, 0x66, 0x77), a=PP_ALIGN.CENTER)

    # Save
    prs.save(output_path)
    return output_path


# ============================================================
# MAIN
# ============================================================
def main():
    if len(sys.argv) < 2:
        # If no argument, try to find the most recent spreadsheet
        folder = os.path.join(os.path.expanduser('~'), 'Desktop', 'ATENDIMENTOS', 'RCPCC - Atendimentos')
        if os.path.isdir(folder):
            xlsx_files = [
                os.path.join(folder, f)
                for f in os.listdir(folder)
                if f.startswith('Tabela atendimento') and f.endswith('.xlsx')
            ]
            if xlsx_files:
                # Sort by modification time, most recent first
                xlsx_files.sort(key=os.path.getmtime, reverse=True)
                print("Planilhas disponíveis (mais recentes primeiro):")
                for i, f in enumerate(xlsx_files[:15]):
                    name = os.path.basename(f)
                    print(f"  [{i + 1}] {name}")
                print()
                choice = input("Escolha o número da planilha (ou Enter para a mais recente): ").strip()
                if choice == '':
                    xlsx_path = xlsx_files[0]
                elif choice.isdigit() and 1 <= int(choice) <= len(xlsx_files):
                    xlsx_path = xlsx_files[int(choice) - 1]
                else:
                    print("Opção inválida.")
                    sys.exit(1)
            else:
                print("Nenhuma planilha encontrada na pasta RCPCC.")
                sys.exit(1)
        else:
            print("Uso: python rcpcc_generator.py <caminho_planilha.xlsx>")
            sys.exit(1)
    else:
        xlsx_path = sys.argv[1]

    if not os.path.isfile(xlsx_path):
        print(f"ERRO: Arquivo não encontrado: {xlsx_path}")
        sys.exit(1)

    print(f"📊 Lendo planilha: {os.path.basename(xlsx_path)}")
    data = read_spreadsheet(xlsx_path)

    print(f"👤 Cliente: {data['client_name']}")
    print(f"💳 Dívidas: {data['total_dividas']} ({data['bancos_unicos']} bancos)")
    print(f"💰 Total: {fmt_brl(data['total_geral'])}")
    print(f"✅ Pagamento com desconto: {fmt_brl(data['pagamento'])} ({data['pct_desconto']:.0f}% desconto)")
    print(f"📈 Economia: {fmt_brl(data['economia'])}")
    print()

    # Output path: same folder as spreadsheet
    folder = os.path.dirname(xlsx_path)
    safe_name = re.sub(r'[^\w\s-]', '', data['client_name']).strip()
    output_name = f"Apresentacao_RCPCC_{safe_name}.pptx"
    output_path = os.path.join(folder, output_name)

    print(f"🎨 Gerando apresentação...")
    generate_presentation(data, output_path)
    print(f"✅ Apresentação salva em: {output_path}")
    print()
    print("Pronto! Abra o arquivo .pptx no PowerPoint.")


if __name__ == '__main__':
    main()
